from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── In-memory store ──────────────────────────────────────────
presentations = {}

# ── 8 colour themes ──────────────────────────────────────────
THEMES = [
    {"bg":"0D1B2A","accent":"E94560","heading":"FFFFFF","body":"CBD5E1","card":"1E2D3D","sub":"94A3B8"},
    {"bg":"1A1A2E","accent":"7C3AED","heading":"FFFFFF","body":"C4B5FD","card":"16213E","sub":"8B5CF6"},
    {"bg":"064E3B","accent":"10B981","heading":"FFFFFF","body":"A7F3D0","card":"065F46","sub":"6EE7B7"},
    {"bg":"1E1B4B","accent":"F59E0B","heading":"FFFFFF","body":"FDE68A","card":"312E81","sub":"FCD34D"},
    {"bg":"1F2937","accent":"06B6D4","heading":"FFFFFF","body":"A5F3FC","card":"111827","sub":"67E8F9"},
    {"bg":"450A0A","accent":"F87171","heading":"FFFFFF","body":"FECACA","card":"7F1D1D","sub":"FCA5A5"},
    {"bg":"0C4A6E","accent":"38BDF8","heading":"FFFFFF","body":"BAE6FD","card":"075985","sub":"7DD3FC"},
    {"bg":"1A2F1A","accent":"84CC16","heading":"FFFFFF","body":"D9F99D","card":"14532D","sub":"BEF264"},
]

# ── Helpers ───────────────────────────────────────────────────
def _rgb(h):
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def _theme(title):
    return THEMES[hash(title) % len(THEMES)]

def _bg(slide, hex_color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(hex_color)

def _box(slide, x, y, w, h, hex_color):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = _rgb(hex_color)
    s.line.fill.background()
    return s

def _oval(slide, x, y, w, h, hex_color):
    s = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = _rgb(hex_color)
    s.line.fill.background()
    return s

def _txt(slide, text, x, y, w, h, size, hex_color, bold=False, italic=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = str(text)
    r.font.size = Pt(size)
    r.font.color.rgb = _rgb(hex_color)
    r.font.bold = bold
    r.font.italic = italic


# ══════════════════════════════════════════════════════════════
#  CREATE PRESENTATION  — title slide
# ══════════════════════════════════════════════════════════════
def create_presentation(title, subtitle="", filename="output.pptx"):
    if not filename.endswith(".pptx"):
        filename += ".pptx"

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    t = _theme(title)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, t["bg"])

    # decorative elements
    _box(slide,  0,    0,    0.20, 7.50, t["accent"])   # left bar
    _oval(slide, 10.5, -1.2, 4.00, 4.00, t["accent"])   # big circle
    _oval(slide, 11.8,  4.8, 2.20, 2.20, t["card"])     # small circle

    # content card
    _box(slide, 0.5, 1.9, 8.8, 3.8, t["card"])

    # title + underline + subtitle
    _txt(slide, title.upper(), 0.75, 2.1, 8.2, 1.8, 38, t["heading"], bold=True)
    _box(slide, 0.75, 3.9, 2.8, 0.07, t["accent"])
    _txt(slide, subtitle or "A Comprehensive AI-Generated Presentation",
         0.75, 4.1, 8.2, 0.9, 18, t["sub"], italic=True)

    # footer
    _box(slide, 0, 6.95, 13.33, 0.55, t["accent"])
    _txt(slide, "Created with Presentation LLM  •  Powered by Mistral + MCP Tools",
         0.4, 6.97, 12.5, 0.45, 11, t["bg"])

    presentations[filename] = {"prs": prs, "theme": t}
    prs.save(filename)
    return f"✅ Presentation '{title}' created → {filename}"


# ══════════════════════════════════════════════════════════════
#  ADD SLIDE  — blank styled canvas
# ══════════════════════════════════════════════════════════════
def add_slide(filename="output.pptx", layout=1):
    if not filename.endswith(".pptx"):
        filename += ".pptx"

    if filename not in presentations:
        if not os.path.exists(filename):
            return f"❌ '{filename}' not found. Create it first."
        presentations[filename] = {"prs": Presentation(filename), "theme": THEMES[0]}

    prs = presentations[filename]["prs"]
    t   = presentations[filename]["theme"]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, t["bg"])
    _box(slide, 0,    0,     13.33, 1.30, t["card"])    # header
    _box(slide, 0,    0,     0.20,  7.50, t["accent"])  # left bar
    _box(slide, 0,    7.10,  13.33, 0.40, t["accent"])  # footer

    prs.save(filename)
    return f"✅ Slide added → Total: {len(prs.slides)} slides"


# ══════════════════════════════════════════════════════════════
#  WRITE TEXT  — 2-column numbered bullet cards
# ══════════════════════════════════════════════════════════════
def write_text(filename="output.pptx", slide_index=1, title="", body=""):
    if not filename.endswith(".pptx"):
        filename += ".pptx"

    if filename not in presentations:
        if not os.path.exists(filename):
            return f"❌ '{filename}' not found. Create it first."
        presentations[filename] = {"prs": Presentation(filename), "theme": THEMES[0]}

    prs = presentations[filename]["prs"]
    t   = presentations[filename]["theme"]

    if slide_index < 1 or slide_index > len(prs.slides):
        return f"❌ Invalid slide index {slide_index}."

    slide = prs.slides[slide_index - 1]

    # ── Title ──
    if title:
        _txt(slide, str(title), 0.35, 0.08, 12.0, 1.1, 26, t["heading"], bold=True)
        _box(slide, 0.35, 1.1, 2.0, 0.06, t["accent"])

    # ── Bullet cards ──
    if body:
        lines = [l.strip() for l in body.replace("\\n","\n").split("\n") if l.strip()][:6]
        colors = [t["accent"], t["sub"]]
        CW, CH = 6.0, 1.42
        GX, GY = 0.28, 0.16
        SX, SY = 0.35, 1.25

        for i, line in enumerate(lines):
            col = i % 2
            row = i // 2
            cx  = SX + col * (CW + GX)
            cy  = SY + row * (CH + GY)
            ac  = colors[i % 2]

            # card background with border
            card = slide.shapes.add_shape(1, Inches(cx), Inches(cy), Inches(CW), Inches(CH))
            card.fill.solid()
            card.fill.fore_color.rgb = _rgb(t["card"])
            card.line.color.rgb      = _rgb(ac)
            card.line.width          = Pt(1.0)

            # left colour strip
            _box(slide, cx, cy, 0.09, CH, ac)

            # number circle
            _oval(slide, cx+0.14, cy+0.10, 0.40, 0.40, ac)
            _txt(slide, str(i+1), cx+0.16, cy+0.10, 0.36, 0.40,
                 13, t["bg"], bold=True, align=PP_ALIGN.CENTER)

            # bullet text
            _txt(slide, str(line), cx+0.62, cy+0.08, CW-0.74, CH-0.16,
                 12, t["body"])

    prs.save(filename)
    return f"✅ Content written to slide {slide_index}."


# ══════════════════════════════════════════════════════════════
#  SEARCH WEB  — local knowledge base
# ══════════════════════════════════════════════════════════════
def search_web(topic=""):
    kb = {
        "python":                  "Python is a high-level language used widely in data science, web development, and AI.",
        "ai":                      "Artificial Intelligence enables machines to learn, reason, and solve problems.",
        "artificial intelligence": "AI is transforming industries via automation, NLP, computer vision, and generative models.",
        "machine learning":        "ML lets systems learn from data and improve automatically without explicit programming.",
        "deep learning":           "Deep Learning uses multi-layer neural networks for complex tasks like image and speech recognition.",
        "blockchain":              "Blockchain is a decentralised ledger ensuring secure, transparent, tamper-proof transactions.",
        "data science":            "Data Science extracts insights from data using statistics, programming, and domain expertise.",
        "cybersecurity":           "Cybersecurity protects systems and networks from digital attacks and unauthorised access.",
        "cloud computing":         "Cloud computing delivers servers, storage, and software over the internet on demand.",
        "climate change":          "Climate change refers to long-term shifts in global temperatures driven by human activity.",
        "renewable energy":        "Renewable energy from solar, wind, and hydro offers sustainable alternatives to fossil fuels.",
        "space exploration":       "Space exploration expands scientific knowledge and drives breakthrough technologies.",
        "quantum computing":       "Quantum computing uses quantum mechanics to solve problems faster than classical computers.",
        "iot":                     "IoT connects billions of physical devices to the internet for smart automation.",
        "5g":                      "5G delivers ultra-fast wireless speeds and low latency for next-gen connected devices.",
    }
    tl = topic.lower().strip()
    for key, val in kb.items():
        if key in tl or tl in key:
            return f"✅ {val}"
    return f"ℹ️ No data found for '{topic}'."
